import os
import sqlite3
import base64
import requests
from datetime import datetime
from uuid import uuid4
from urllib.parse import quote_plus

from fastapi import FastAPI, Request, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse

from slowapi import Limiter
from slowapi.util import get_remote_address

from groq import Groq
from fpdf import FPDF
from pptx import Presentation
import replicate

# ==========================================================
# CONFIG
# ==========================================================
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
REPLICATE_API_TOKEN = os.getenv("REPLICATE_API_TOKEN")
BASE_URL = os.getenv("BASE_URL", "https://levo-server.onrender.com")

client = Groq(api_key=GROQ_API_KEY)
replicate_client = replicate.Client(api_token=REPLICATE_API_TOKEN)

app = FastAPI()
limiter = Limiter(key_func=get_remote_address)
app.state.limiter = limiter

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

os.makedirs("files", exist_ok=True)

# ==========================================================
# SYSTEM PROMPT
# ==========================================================
SYSTEM_PROMPT = """
You are lEvO, a calm and accurate AI assistant.
No fake links. No hallucinations.
"""

# ==========================================================
# MEMORY
# ==========================================================
conn = sqlite3.connect("memory.db", check_same_thread=False)
cursor = conn.cursor()
cursor.execute("CREATE TABLE IF NOT EXISTS memory (user_id TEXT, text TEXT)")
conn.commit()

def get_memory(uid):
    cursor.execute("SELECT text FROM memory WHERE user_id=?", (uid,))
    return "\n".join(row[0] for row in cursor.fetchall())

def save_memory(uid, msg):
    cursor.execute("INSERT INTO memory VALUES (?, ?)", (uid, msg))
    conn.commit()

# ==========================================================
# LINK SEARCH
# ==========================================================
def get_links(query, site=None):
    q = f"{query} site:{site}" if site else query
    url = f"https://ddg-api.herokuapp.com/search?query={quote_plus(q)}"
    data = requests.get(url).json()

    links = []
    for r in data.get("results", [])[:5]:
        links.append(f"{r['title']} → {r['url']}")

    return "\n".join(links)

# ==========================================================
# CHAT
# ==========================================================
@app.post("/chat")
@limiter.limit("15/minute")
async def chat(request: Request):
    data = await request.json()
    prompt = data.get("prompt", "")
    user_id = data.get("user_id", "guest")

    memory = get_memory(user_id)
    today = datetime.now().strftime("%d %B %Y")

    wants_amazon = "amazon" in prompt.lower()
    wants_youtube = "youtube" in prompt.lower()

    links = ""
    if wants_amazon:
        links = get_links(prompt, "amazon.com")
    elif wants_youtube:
        links = get_links(prompt, "youtube.com")

    messages = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "system", "content": f"Date: {today}"},
        {"role": "system", "content": f"Memory:\n{memory}"},
        {"role": "user", "content": prompt}
    ]

    if links:
        messages.insert(-1, {"role": "system", "content": f"Verified links:\n{links}"})

    completion = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=messages,
        max_tokens=2048
    )

    reply = completion.choices[0].message.content
    save_memory(user_id, prompt)
    save_memory(user_id, reply)

    return {"response": reply}

# ==========================================================
# IMAGE ANALYSIS (FIXED)
# ==========================================================
@app.post("/vision")
async def vision(file: UploadFile = File(...)):
    image_bytes = await file.read()
    image_b64 = base64.b64encode(image_bytes).decode()

    output = replicate_client.run(
        "salesforce/blip",
        input={"image": f"data:image/jpeg;base64,{image_b64}"}
    )

    return {"response": output}

# ==========================================================
# IMAGE GENERATION (FIXED)
# ==========================================================
@app.post("/image-generate")
def image_generate(prompt: str):
    output = replicate_client.run(
        "stability-ai/sdxl:39ed52f2a78e934b3ba6e2a89f5b1c712de7f50c08b2c2a5f5c99f8b6f10c6fb",
        input={"prompt": prompt}
    )

    return {
        "image": output[0],
        "download": output[0]
    }

# ==========================================================
# PDF
# ==========================================================
@app.post("/create-pdf")
def create_pdf(text: str):
    path = f"files/{uuid4()}.pdf"
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 8, text)
    pdf.output(path)

    return {"file": f"{BASE_URL}/{path}", "type": "pdf"}

# ==========================================================
# PPT
# ==========================================================
@app.post("/create-ppt")
def create_ppt(text: str):
    path = f"files/{uuid4()}.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Generated by lEvO"
    slide.placeholders[1].text = text
    prs.save(path)

    return {"file": f"{BASE_URL}/{path}", "type": "pptx"}

# ==========================================================
# FILE SERVE
# ==========================================================
@app.get("/files/{filename}")
def serve_file(filename: str):
    return FileResponse(f"files/{filename}")
